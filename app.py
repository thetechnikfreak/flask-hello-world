from flask import Flask, render_template, request, redirect, send_file, url_for
from datetime import datetime
import feedparser
from pptx import Presentation
from pptx.util import Inches
import requests
from io import BytesIO
import tempfile
import os
import uuid

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'static/uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# In-Memory-Event-Liste (später ggf. Datei oder DB)
events = []

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/feed")
def feed_view():
    return render_template("feed.html", events=events)

@app.route("/submit", methods=["GET", "POST"])
def submit():
    if request.method == "POST":
        title = request.form.get("title", "")[:30]
        description = request.form.get("description", "")[:143]
        date = request.form.get("date")
        file = request.files.get("image")

        if file:
            filename = f"{uuid.uuid4().hex}.jpg"
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            image_url = url_for('static', filename=f'uploads/{filename}', _external=True)
        else:
            image_url = None

        events.append({
            "title": title,
            "description": description,
            "date": date,
            "image": image_url
        })

        return redirect("/feed")

    return render_template("form.html")

@app.route("/generate-pptx")
def generate_pptx():
    prs = Presentation()

    for event in events:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = event["title"]
        desc = event["description"]

        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1.5))
        tf = txBox.text_frame
        tf.text = title + "\n" + desc

        img_url = event.get("image")
        if img_url:
            try:
                img_data = requests.get(img_url).content
                img_stream = BytesIO(img_data)
                slide.shapes.add_picture(img_stream, Inches(3), Inches(2), Inches(4), Inches(4))
            except Exception as e:
                print(f"Fehler beim Bildabruf: {e}")

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return send_file(tmp.name, as_attachment=True, download_name="mamoba-veranstaltungen.pptx")

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port)
